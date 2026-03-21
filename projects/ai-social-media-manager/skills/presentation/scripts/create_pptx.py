#!/usr/bin/env python3
"""
create_pptx.py -- Build branded presentations from a JSON slide manifest.

Usage:
    python3 create_pptx.py --input slides.json --output deck.pptx
    python3 create_pptx.py --input slides.json --output deck.pptx --context conference --accent gold

Arguments:
    --input     Path to JSON slide manifest (or - for stdin)
    --output    Output .pptx file path
    --context   Presentation context for font sizing: large_venue, conference, meeting, screen_share, pdf
    --accent    Accent color: gold, sage, blush, lavender (default: gold)
"""

import argparse
import json
import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Brand Constants
# ---------------------------------------------------------------------------

COLORS = {
    "ivory": RGBColor(0xF7, 0xF4, 0xEF),
    "warm_linen": RGBColor(0xF0, 0xEC, 0xE4),
    "charcoal": RGBColor(0x3A, 0x35, 0x2E),
    "gold": RGBColor(0xB8, 0xA0, 0x6A),
    "light_gold": RGBColor(0xD4, 0xC4, 0x8E),
    "sage": RGBColor(0x8F, 0xAB, 0x8A),
    "blush": RGBColor(0xD4, 0xB0, 0xA8),
    "lavender": RGBColor(0xC4, 0xB8, 0xCC),
    "text_secondary": RGBColor(0x7A, 0x72, 0x68),
    "text_muted": RGBColor(0xB0, 0xA8, 0x98),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

ACCENT_MAP = {
    "gold": COLORS["gold"],
    "sage": COLORS["sage"],
    "blush": COLORS["blush"],
    "lavender": COLORS["lavender"],
}

# Font families with fallbacks (python-pptx embeds the name; the viewer
# resolves it). If Playfair Display is not installed on the viewing machine,
# PowerPoint will substitute a similar serif font automatically.
FONT_HEADLINE = "Playfair Display"
FONT_BODY = "DM Sans"
FONT_HEADLINE_FALLBACK = "Georgia"
FONT_BODY_FALLBACK = "Calibri"

FONT_SIZES = {
    "large_venue":  {"title": Pt(44), "subtitle": Pt(24), "body": Pt(32), "caption": Pt(18), "number": Pt(96)},
    "conference":   {"title": Pt(40), "subtitle": Pt(22), "body": Pt(28), "caption": Pt(16), "number": Pt(88)},
    "meeting":      {"title": Pt(36), "subtitle": Pt(20), "body": Pt(24), "caption": Pt(14), "number": Pt(80)},
    "screen_share": {"title": Pt(32), "subtitle": Pt(20), "body": Pt(24), "caption": Pt(16), "number": Pt(72)},
    "pdf":          {"title": Pt(28), "subtitle": Pt(18), "body": Pt(20), "caption": Pt(14), "number": Pt(64)},
}

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_WIDTH = SLIDE_WIDTH - (2 * MARGIN)
TITLE_TOP = Inches(0.75)
TITLE_HEIGHT = Inches(1.2)
CONTENT_TOP = Inches(2.2)
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=None):
    """Set slide background to a solid color (default: ivory)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or COLORS["ivory"]


def add_text_box(slide, left, top, width, height, text, font_name, font_size,
                 font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box with styled text. Returns the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color or COLORS["charcoal"]
    run.font.bold = bold

    return txBox


def add_accent_bar(slide, accent_color, left=None, top=None, width=None, height=None):
    """Add a thin accent color bar (decorative element)."""
    left = left or MARGIN
    top = top or Inches(2.0)
    width = width or Inches(2.0)
    height = height or Inches(0.06)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()
    return shape


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, slide_data, sizes, accent):
    """Title slide: large title, subtitle, accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide)

    title = slide_data.get("title", "Untitled Presentation")
    subtitle = slide_data.get("subtitle", "")

    # Title centered vertically
    add_text_box(
        slide, MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(1.5),
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER
    )

    # Accent bar below title
    bar_width = Inches(3.0)
    bar_left = (SLIDE_WIDTH - bar_width) // 2
    add_accent_bar(slide, accent, left=bar_left, top=Inches(4.2), width=bar_width)

    # Subtitle
    if subtitle:
        add_text_box(
            slide, MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(0.8),
            subtitle, FONT_BODY, sizes["subtitle"],
            font_color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_agenda_slide(prs, slide_data, sizes, accent):
    """Agenda/TOC slide: numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "Agenda")
    items = slide_data.get("items", [])

    # Title
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    # Items
    for i, item in enumerate(items[:7]):
        item_top = CONTENT_TOP + Inches(i * 0.7)
        # Number
        add_text_box(
            slide, MARGIN, item_top, Inches(0.6), Inches(0.6),
            f"{i + 1}.", FONT_HEADLINE, sizes["body"],
            font_color=accent, bold=True
        )
        # Text
        add_text_box(
            slide, Inches(1.5), item_top, CONTENT_WIDTH - Inches(0.75), Inches(0.6),
            item, FONT_BODY, sizes["body"],
            font_color=COLORS["charcoal"]
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_section_divider(prs, slide_data, sizes, accent):
    """Section divider: large section title centered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["warm_linen"])

    title = slide_data.get("title", "")

    add_text_box(
        slide, MARGIN, Inches(2.8), CONTENT_WIDTH, Inches(1.5),
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER
    )

    bar_width = Inches(2.0)
    bar_left = (SLIDE_WIDTH - bar_width) // 2
    add_accent_bar(slide, accent, left=bar_left, top=Inches(4.5), width=bar_width)

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_content_slide(prs, slide_data, sizes, accent):
    """Content slide: assertion headline + body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "")
    body = slide_data.get("body", "")
    image_path = slide_data.get("image")

    # Title (assertion headline)
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    if image_path and os.path.exists(image_path):
        # Two-column: text left, image right
        text_width = CONTENT_WIDTH * 0.55
        img_left = MARGIN + text_width + Inches(0.3)
        img_width = CONTENT_WIDTH * 0.4

        if body:
            add_text_box(
                slide, MARGIN, CONTENT_TOP, text_width, CONTENT_HEIGHT,
                body, FONT_BODY, sizes["body"],
                font_color=COLORS["text_secondary"]
            )

        try:
            slide.shapes.add_picture(
                image_path, img_left, CONTENT_TOP, img_width
            )
        except Exception as e:
            print(f"Warning: Could not insert image {image_path}: {e}", file=sys.stderr)
    else:
        # Full-width body
        if body:
            add_text_box(
                slide, MARGIN, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT,
                body, FONT_BODY, sizes["body"],
                font_color=COLORS["text_secondary"]
            )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_data_slide(prs, slide_data, sizes, accent):
    """Data slide: assertion headline + image/chart placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "")
    image_path = slide_data.get("image")
    caption = slide_data.get("caption", "")

    # Assertion headline
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    # Image or placeholder
    img_area_left = Inches(2.0)
    img_area_width = SLIDE_WIDTH - Inches(4.0)
    img_area_top = CONTENT_TOP
    img_area_height = Inches(4.0)

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(
                image_path, img_area_left, img_area_top, img_area_width
            )
        except Exception as e:
            print(f"Warning: Could not insert image {image_path}: {e}", file=sys.stderr)
    else:
        # Placeholder rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, img_area_left, img_area_top,
            img_area_width, img_area_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["warm_linen"]
        shape.line.color.rgb = COLORS["text_muted"]
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "[Insert chart or data visualization]"
        run.font.name = FONT_BODY
        run.font.size = sizes["caption"]
        run.font.color.rgb = COLORS["text_muted"]

    # Caption
    if caption:
        add_text_box(
            slide, MARGIN, Inches(6.3), CONTENT_WIDTH, Inches(0.5),
            caption, FONT_BODY, sizes["caption"],
            font_color=COLORS["text_muted"], alignment=PP_ALIGN.CENTER
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_comparison_slide(prs, slide_data, sizes, accent):
    """Comparison slide: 2-3 columns with parallel content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "")
    columns = slide_data.get("columns", [])

    # Title
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    # Columns
    num_cols = min(len(columns), 3)
    if num_cols == 0:
        add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    col_gap = Inches(0.4)
    total_gap = col_gap * (num_cols - 1)
    col_width = (CONTENT_WIDTH - total_gap) / num_cols

    for i, col in enumerate(columns[:3]):
        col_left = MARGIN + (col_width + col_gap) * i
        col_title = col.get("title", "")
        col_body = col.get("body", "")

        # Column header (with accent background)
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col_left, CONTENT_TOP,
            col_width, Inches(0.7)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = accent
        header_shape.line.fill.background()

        tf = header_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = col_title
        run.font.name = FONT_BODY
        run.font.size = sizes["body"]
        run.font.color.rgb = COLORS["white"]
        run.font.bold = True

        # Column body
        if col_body:
            add_text_box(
                slide, col_left, CONTENT_TOP + Inches(0.9), col_width, Inches(3.5),
                col_body, FONT_BODY, sizes["body"],
                font_color=COLORS["text_secondary"]
            )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_quote_slide(prs, slide_data, sizes, accent):
    """Quote slide: large quote text + attribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["warm_linen"])

    quote = slide_data.get("quote", "")
    attribution = slide_data.get("attribution", "")

    # Opening quotation mark (decorative)
    add_text_box(
        slide, Inches(1.5), Inches(1.5), Inches(1.0), Inches(1.0),
        "\u201C", FONT_HEADLINE, Pt(120),
        font_color=accent, alignment=PP_ALIGN.LEFT
    )

    # Quote text
    add_text_box(
        slide, Inches(2.0), Inches(2.5), SLIDE_WIDTH - Inches(4.0), Inches(3.0),
        quote, FONT_HEADLINE, sizes["body"],
        font_color=COLORS["charcoal"], alignment=PP_ALIGN.LEFT
    )

    # Attribution
    if attribution:
        add_text_box(
            slide, Inches(2.0), Inches(5.8), SLIDE_WIDTH - Inches(4.0), Inches(0.6),
            f"\u2014 {attribution}", FONT_BODY, sizes["caption"],
            font_color=COLORS["text_secondary"], alignment=PP_ALIGN.LEFT
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_big_number_slide(prs, slide_data, sizes, accent):
    """Big number slide: one large stat + context sentence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    number = slide_data.get("number", "0")
    context = slide_data.get("context", "")

    # Large number
    add_text_box(
        slide, MARGIN, Inches(2.0), CONTENT_WIDTH, Inches(2.5),
        str(number), FONT_HEADLINE, sizes["number"],
        font_color=accent, bold=True, alignment=PP_ALIGN.CENTER
    )

    # Context sentence
    if context:
        add_text_box(
            slide, Inches(2.0), Inches(4.8), SLIDE_WIDTH - Inches(4.0), Inches(1.0),
            context, FONT_BODY, sizes["body"],
            font_color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_process_slide(prs, slide_data, sizes, accent):
    """Process slide: 3-5 numbered steps in a horizontal flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "")
    steps = slide_data.get("steps", [])

    # Title
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    # Steps in horizontal layout
    num_steps = min(len(steps), 5)
    if num_steps == 0:
        add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    step_gap = Inches(0.3)
    total_gap = step_gap * (num_steps - 1)
    step_width = (CONTENT_WIDTH - total_gap) / num_steps
    circle_size = Inches(0.7)

    for i, step in enumerate(steps[:5]):
        step_left = MARGIN + (step_width + step_gap) * i
        center_x = step_left + (step_width - circle_size) / 2

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, center_x, CONTENT_TOP, circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = FONT_BODY
        run.font.size = sizes["body"]
        run.font.color.rgb = COLORS["white"]
        run.font.bold = True

        # Arrow between steps (except after last)
        if i < num_steps - 1:
            arrow_left = step_left + step_width + Inches(0.02)
            arrow_top = CONTENT_TOP + circle_size / 2 - Inches(0.05)
            add_text_box(
                slide, arrow_left, arrow_top, step_gap - Inches(0.04), Inches(0.3),
                "\u2192", FONT_BODY, sizes["caption"],
                font_color=COLORS["text_muted"], alignment=PP_ALIGN.CENTER
            )

        # Step text
        step_text = step if isinstance(step, str) else step.get("text", "")
        add_text_box(
            slide, step_left, CONTENT_TOP + Inches(1.0), step_width, Inches(2.5),
            step_text, FONT_BODY, sizes["body"],
            font_color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_summary_slide(prs, slide_data, sizes, accent):
    """Summary slide: 3-5 key takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "Key Takeaways")
    points = slide_data.get("points", [])

    # Title
    add_text_box(
        slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True
    )

    add_accent_bar(slide, accent)

    # Points with accent markers
    for i, point in enumerate(points[:5]):
        point_top = CONTENT_TOP + Inches(i * 0.85)

        # Accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN, point_top + Inches(0.12),
            Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        # Point text
        add_text_box(
            slide, Inches(1.2), point_top, CONTENT_WIDTH - Inches(0.45), Inches(0.7),
            point, FONT_BODY, sizes["body"],
            font_color=COLORS["charcoal"]
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


def build_cta_slide(prs, slide_data, sizes, accent):
    """CTA / Closing slide: call to action + contact info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    title = slide_data.get("title", "Next Steps")
    action = slide_data.get("action", "")
    contact = slide_data.get("contact", "")

    # Title
    add_text_box(
        slide, MARGIN, Inches(2.0), CONTENT_WIDTH, Inches(1.5),
        title, FONT_HEADLINE, sizes["title"],
        font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER
    )

    # Accent bar
    bar_width = Inches(3.0)
    bar_left = (SLIDE_WIDTH - bar_width) // 2
    add_accent_bar(slide, accent, left=bar_left, top=Inches(3.7), width=bar_width)

    # CTA action
    if action:
        add_text_box(
            slide, Inches(2.0), Inches(4.0), SLIDE_WIDTH - Inches(4.0), Inches(1.0),
            action, FONT_BODY, sizes["body"],
            font_color=accent, bold=True, alignment=PP_ALIGN.CENTER
        )

    # Contact info
    if contact:
        add_text_box(
            slide, Inches(2.0), Inches(5.5), SLIDE_WIDTH - Inches(4.0), Inches(0.8),
            contact, FONT_BODY, sizes["caption"],
            font_color=COLORS["text_secondary"], alignment=PP_ALIGN.CENTER
        )

    add_speaker_notes(slide, slide_data.get("notes", ""))
    return slide


# ---------------------------------------------------------------------------
# Builder dispatch
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    "title": build_title_slide,
    "agenda": build_agenda_slide,
    "section-divider": build_section_divider,
    "section_divider": build_section_divider,
    "content": build_content_slide,
    "data": build_data_slide,
    "chart": build_data_slide,
    "comparison": build_comparison_slide,
    "quote": build_quote_slide,
    "big-number": build_big_number_slide,
    "big_number": build_big_number_slide,
    "stat": build_big_number_slide,
    "process": build_process_slide,
    "flow": build_process_slide,
    "summary": build_summary_slide,
    "recap": build_summary_slide,
    "cta": build_cta_slide,
    "closing": build_cta_slide,
}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def create_presentation(manifest, context="conference", accent_name="gold"):
    """Build a presentation from a slide manifest dict."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    sizes = FONT_SIZES.get(context, FONT_SIZES["conference"])
    accent = ACCENT_MAP.get(accent_name, COLORS["gold"])

    slides_data = manifest.get("slides", [])

    MAX_SLIDES = 100
    if len(slides_data) > MAX_SLIDES:
        print(f"Warning: Slide count ({len(slides_data)}) exceeds {MAX_SLIDES}. Truncating.", file=sys.stderr)
        slides_data = slides_data[:MAX_SLIDES]

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content").lower().strip()
        builder = SLIDE_BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data, sizes, accent)

    return prs


def main():
    parser = argparse.ArgumentParser(description="Build branded PPTX from JSON manifest")
    parser.add_argument("--input", "-i", required=True, help="JSON manifest path (or - for stdin)")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx path")
    parser.add_argument("--context", "-c", default="conference",
                        choices=["large_venue", "conference", "meeting", "screen_share", "pdf"],
                        help="Presentation context for font sizing")
    parser.add_argument("--accent", "-a", default="gold",
                        choices=["gold", "sage", "blush", "lavender"],
                        help="Accent color (content pillar)")
    args = parser.parse_args()

    # Load manifest
    if args.input == "-":
        manifest = json.load(sys.stdin)
    else:
        with open(args.input, "r") as f:
            manifest = json.load(f)

    # Override context/accent from manifest if present
    context = manifest.get("context", args.context)
    accent = manifest.get("accent", args.accent)

    prs = create_presentation(manifest, context=context, accent_name=accent)

    # Ensure output directory exists
    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(out_path))
    print(f"Presentation saved: {out_path}")
    print(f"  Slides: {len(manifest.get('slides', []))}")
    print(f"  Context: {context}")
    print(f"  Accent: {accent}")


if __name__ == "__main__":
    main()
